"""Build the MPG executive presentation.

Every figure here was computed against the database with the same logic the semantic
model uses, at the grain the report displays, so the deck and the tool cannot
disagree. The palette is the report's own: a reader who opens the .pbix after the
meeting should feel they are looking at the same asset.
"""
import pathlib

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

OUT = (pathlib.Path(__file__).resolve().parents[1]
       / "MPG_PeopleAnalytics_Executive.pptx")

NAVY, BLUE, LIGHT, PALE = "001552", "3E72C4", "9CBCE8", "C9DAF3"
GREY, WHITE, BODY, MUTED = "F2F4F7", "FFFFFF", "475467", "667085"
OCHRE, BORDER = "96691A", "E4E7EC"
HEAD_FONT, BODY_FONT = "Cambria", "Calibri"

W, H = 13.333, 7.5
prs = Presentation()
prs.slide_width, prs.slide_height = Inches(W), Inches(H)
BLANK = prs.slide_layouts[6]


def rgb(h):
    return RGBColor.from_string(h)


def slide(dark=False):
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = rgb(NAVY if dark else GREY)
    bg.line.fill.background()
    bg.shadow.inherit = False
    return s


def card(s, x, y, w, h, fill=WHITE, line=BORDER):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
                             Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(fill)
    if line:
        shp.line.color.rgb = rgb(line)
        shp.line.width = Pt(0.75)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    shp.adjustments[0] = 0.06
    return shp


def text(s, x, y, w, h, blocks, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """blocks: list of (string, size, color, bold, font, space_after_pt)."""
    box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    for i, (txt, size, color, bold, font, after) in enumerate(blocks):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(after)
        run = p.add_run()
        run.text = txt
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = rgb(color)
        run.font.name = font
    return box


def link(s, x, y, w, h, label, url, size, colour):
    """A clickable line of text. A URL nobody can click is a URL nobody will type."""
    box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    run = tf.paragraphs[0].add_run()
    run.text = label
    run.font.size = Pt(size)
    run.font.name = BODY_FONT
    run.hyperlink.address = url
    # The hyperlink would otherwise inherit the theme's link colour, which is not ours.
    run.font.color.rgb = rgb(colour)
    return box


def title(s, t, sub=None, dark=False):
    text(s, 0.75, 0.52, W - 1.5, 1.0,
         [(t, 32, WHITE if dark else NAVY, True, HEAD_FONT, 0)])
    if sub:
        text(s, 0.75, 1.42, W - 1.5, 0.55,
             [(sub, 14, PALE if dark else MUTED, False, BODY_FONT, 0)])


def stat(s, x, y, w, label, value, sub, color=NAVY, h=1.75):
    card(s, x, y, w, h)
    text(s, x + 0.28, y + 0.24, w - 0.56, 0.42,
         [(label, 11, MUTED, False, BODY_FONT, 0)])
    text(s, x + 0.28, y + 0.62, w - 0.56, 0.62,
         [(value, 30, color, False, HEAD_FONT, 0)])
    text(s, x + 0.28, y + 1.26, w - 0.56, 0.42,
         [(sub, 10, MUTED, False, BODY_FONT, 0)])


def bullets(s, x, y, w, h, items, size=13, gap=9):
    blocks = []
    for head, rest in items:
        blocks.append((head, size, NAVY, True, BODY_FONT, 2))
        blocks.append((rest, size - 1, BODY, False, BODY_FONT, gap))
    return text(s, x, y, w, h, blocks)


def numbered(s, x, y, n, head, rest, w=3.6):
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(0.5), Inches(0.5))
    circ.fill.solid()
    circ.fill.fore_color.rgb = rgb(NAVY)
    circ.line.fill.background()
    circ.shadow.inherit = False
    tf = circ.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = str(n)
    r.font.size = Pt(15)
    r.font.bold = True
    r.font.color.rgb = rgb(WHITE)
    r.font.name = HEAD_FONT
    text(s, x + 0.7, y + 0.02, w, 2.6,
         [(head, 15, NAVY, True, BODY_FONT, 5),
          (rest, 12, BODY, False, BODY_FONT, 0)])


def chart(s, x, y, w, h, cats, series, kind=XL_CHART_TYPE.COLUMN_CLUSTERED,
          colors=(BLUE,), number_format='[$-409]#,##0.0', legend=False):
    data = CategoryChartData()
    data.categories = cats
    for name, values in series:
        data.add_series(name, values, number_format)
    gf = s.shapes.add_chart(kind, Inches(x), Inches(y), Inches(w), Inches(h), data)
    ch = gf.chart
    ch.has_title = False
    ch.has_legend = legend
    if legend:
        ch.legend.position = XL_LEGEND_POSITION.BOTTOM
        ch.legend.include_in_layout = False
        ch.legend.font.size = Pt(10)
        ch.legend.font.name = BODY_FONT
    if len(ch.series) == 1 and len(colors) > 1:
        # One series, several colors: color the points, not the series. Setting the
        # series fill paints every bar the same, which is what happened the first time.
        for idx, point in enumerate(ch.series[0].points):
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = rgb(colors[idx % len(colors)])
    else:
        for idx, plot_series in enumerate(ch.series):
            plot_series.format.fill.solid()
            plot_series.format.fill.fore_color.rgb = rgb(colors[idx % len(colors)])
    plot = ch.plots[0]
    plot.has_data_labels = True
    plot.data_labels.font.size = Pt(9)
    plot.data_labels.font.name = BODY_FONT
    plot.data_labels.font.color.rgb = rgb(BODY)
    plot.data_labels.number_format = number_format
    plot.data_labels.number_format_is_linked = False
    for axis in (ch.category_axis, ch.value_axis):
        axis.has_major_gridlines = False
        axis.tick_labels.font.size = Pt(9)
        axis.tick_labels.font.name = BODY_FONT
        axis.tick_labels.font.color.rgb = rgb(MUTED)
    ch.value_axis.visible = False
    return ch


# ---------------------------------------------------------------- 1 · title
s = slide(dark=True)
text(s, 1.0, 2.35, 11.3, 1.5,
     [("Attrition is not a pay problem everywhere", 40, WHITE, True, HEAD_FONT, 0)])
text(s, 1.0, 3.75, 10.5, 1.0,
     [("A blanket 6 % across LATAM would cost $5.38M and leave the worst-paid "
       "segment where it is. $1.15M, targeted, does more.", 16, PALE, False, BODY_FONT, 0)])
text(s, 1.0, 5.15, 11.3, 0.9,
     [("MPG People & Corporate Services  ·  Executive Committee", 12, LIGHT, False, BODY_FONT, 4),
      ("4,500 employees · 9 cities · 24 months · 746 segments scanned",
       11, "8FA6CC", False, BODY_FONT, 0)])
# The synthetic data is not a caveat to bury. It is the reason the analysis can be
# scored at all: the anomalies were specified before the data was generated, so
# precision and recall mean something here that they could not mean on real records.
text(s, 1.0, 6.02, 8.8, 0.8,
     [("Built on synthetic data, generated to a documented specification. Real "
       "compensation data is confidential — and because the anomalies were designed in "
       "advance, the analysis can be scored rather than asserted.",
       10.5, "8FA6CC", False, BODY_FONT, 0)])
link(s, 1.0, 6.62, 11.3, 0.3, "github.com/Alejandro-Tarazona/01_PeopleAnalytics",
     "https://github.com/Alejandro-Tarazona/01_PeopleAnalytics", 11, LIGHT)
s.notes_slide.notes_text_frame.text = (
    "The recommendation is not to spend less for its own sake. It is that the blanket "
    "proposal spends where the evidence is not.")

# ---------------------------------------------------------------- 2 · the ask
s = slide()
title(s, "What was asked", "HR has proposed a blanket 6 % salary adjustment across "
      "the LATAM hub. Before approving it, the CEO asked three questions.")
for i, (head, rest) in enumerate([
        ("Are we losing people everywhere,\nor in specific places?",
         "Company attrition is 13.5 %. LATAM is at 16.4 %. Neither number says where."),
        ("Is this actually a money problem?",
         "Low pay and high attrition in the same place is a correlation, not a cause. "
         "Three different problems can produce the same symptom."),
        ("What is the smallest number\nthat fixes most of it?",
         "A blanket rate large enough to fix the worst segment overpays everyone else.")]):
    numbered(s, 0.75 + i * 4.0, 2.5, i + 1, head, rest, w=3.15)
card(s, 0.75, 5.85, W - 1.5, 0.95, fill=NAVY, line=None)
text(s, 1.1, 6.12, W - 2.2, 0.5,
     [("This deck answers them in that order. Every figure is reproducible from the "
       "repository.", 13, WHITE, False, BODY_FONT, 0)])

# ---------------------------------------------------------------- 3 · the answer
s = slide()
title(s, "The answer, in one slide")
stat(s, 0.75, 2.15, 3.85, "The proposal on the table", "$5.38M",
     "blanket 6 % across LATAM", OCHRE)
stat(s, 4.83, 2.15, 3.85, "What we recommend instead", "$1.15M",
     "two segments, targeted", NAVY)
stat(s, 8.9, 2.15, 3.68, "Exposure it addresses", "$5.76M", "1 : 5 against the spend", NAVY)
card(s, 0.75, 4.3, W - 1.5, 2.5)
bullets(s, 1.15, 4.58, W - 2.3, 2.12, [
    ("95 % of the blanket money lands on people already paid at or above their band minimum.",
     "It is spread evenly across a problem that is not."),
    ("23 people in LATAM would still be below band minimum after it.",
     "The rate large enough to fix them is not a rate the budget can carry."),
    ("The targeted adjustment pays for itself in under five months.",
     "And still covers three times its cost under the most conservative assumption we tested."),
], size=14, gap=13)

# ---------------------------------------------------------------- 4 · question 1
s = slide()
title(s, "Specific places", "746 segments were scanned. Eleven are large enough to "
      "judge. Four carry a signal.")
for i, (label, value, sub) in enumerate([
        ("Segments scanned", "746", "city × job × level"),
        ("Large enough to judge", "11", "thirty or more people"),
        ("Carrying a signal", "4", "across three different rules"),
        ("Twelve-month exposure", "$83.3M", "599 voluntary exits")]):
    stat(s, 0.75 + i * 3.13, 2.4, 2.88, label, value, sub,
         NAVY if i < 3 else OCHRE)
card(s, 0.75, 4.55, W - 1.5, 2.25)
text(s, 1.15, 4.9, W - 2.3, 1.6,
     [("Why the middle number matters", 15, NAVY, True, BODY_FONT, 6),
      ("735 segments were set aside on size, not on judgement. Below thirty people an "
       "attrition rate is arithmetic on too few lives to act on — a single resignation "
       "moves it by three points. Reporting on them would not be conservative, it would "
       "be noise presented as findings.", 13, BODY, False, BODY_FONT, 0)])

# ---------------------------------------------------------------- 5 · exits vs dollars
s = slide()
title(s, "Heads and dollars point at different cities")
# Reversed: a bar chart draws its first category at the bottom, so the ranking has to
# be handed over upside down to read top-down on the slide.
cities = ["San Francisco", "Mexico City", "São Paulo", "Miami", "London",
          "Bogotá", "New York", "Singapore", "Dublin"]
card(s, 0.75, 2.15, 6.0, 4.0)
text(s, 1.05, 2.4, 5.4, 0.4, [("Exposure, $M", 12, NAVY, True, BODY_FONT, 0)])
chart(s, 0.95, 2.8, 5.6, 3.15, cities,
      [("Exposure", [4.5, 4.9, 5.4, 7.1, 9.2, 9.7, 10.2, 13.4, 13.5])],
      XL_CHART_TYPE.BAR_CLUSTERED, (LIGHT,), '[$-409]#,##0.0')
card(s, 7.05, 2.15, 5.53, 4.0)
text(s, 7.35, 2.4, 5.0, 0.4, [("Voluntary exits", 12, NAVY, True, BODY_FONT, 0)])
chart(s, 7.25, 2.8, 5.13, 3.15, cities,
      [("Exits", [22, 54, 55, 46, 47, 148, 51, 97, 79])],
      XL_CHART_TYPE.BAR_CLUSTERED, (NAVY,), '[$-409]#,##0')
text(s, 0.75, 6.32, W - 1.5, 0.6,
     [("Bogotá lost 148 people, two and a half times more than any other city, and ranks "
       "fourth in exposure. Dublin lost 79 and ranks first: LATAM salaries are lower, so "
       "the same resignation costs less.", 12, BODY, False, BODY_FONT, 0)])

# ---------------------------------------------------------------- 6 · the rule
s = slide()
title(s, "What the rule declines to flag",
      "Three tests, applied identically to all 746 segments, written before the data "
      "was looked at.")
for i, (head, rest) in enumerate([
        ("A materiality floor", "Thirty people. A rule that cannot act below thirty has "
         "no business reporting below thirty."),
        ("A significance test", "The exact Poisson probability of the observed exits at "
         "the company rate. Below 0.05, or it does not fire."),
        ("Three separate rules", "Pay, internal equity, and organizational. Three "
         "problems that look identical in an attrition report.")]):
    card(s, 0.75 + i * 4.2, 2.35, 3.9, 1.85)
    text(s, 1.05 + i * 4.2, 2.6, 3.3, 1.4,
         [(head, 14, NAVY, True, BODY_FONT, 5), (rest, 11.5, BODY, False, BODY_FONT, 0)])
card(s, 0.75, 4.45, W - 1.5, 2.3, fill="FDF6E9", line="E2D3AE")
text(s, 1.15, 4.75, W - 2.3, 1.75,
     [("Two segments, identical numbers, neither fires", 15, OCHRE, True, BODY_FONT, 6),
      ("Bogotá Operations IC3 and São Paulo Sales IC3: 35 people each, 8 voluntary exits "
       "against 4.7 expected, 1.69× the company baseline, p = 0.107. At that population "
       "an elevated reading has close to a one-in-nine chance of being chance alone.",
       13, BODY, False, BODY_FONT, 7),
      ("Under the “1.5× the baseline” rule of thumb this replaced, both would "
       "have been recommended for a raise.", 13, NAVY, True, BODY_FONT, 0)])

# ------------------------------------------------------- 7-9 · the three findings
FINDINGS = [
    ("Bogotá · Data & Analytics IC3–IC4", "PAY ADJUSTMENT", NAVY,
     [("Compa-ratio", "0.85"), ("Attrition vs baseline", "2.42×"),
      ("Significance", "p < 0.0002"), ("Cost to lift to 0.95", "$811K"),
      ("Exposure carried", "$3.43M"), ("Payback", "4.8 months")],
     "The one segment where money is both the cause and the fix. 150 people paid at "
     "85 % of market midpoint, losing a third of themselves a year. The pay gap is not "
     "in a tail: only 15 % sit below band minimum, because the band floor is 80 % of "
     "midpoint. The whole distribution has drifted."),
    ("Dublin · Engineering IC5", "INTERNAL EQUITY", BLUE,
     [("Compa-ratio", "0.99"), ("Attrition vs baseline", "0.74×"),
      ("Incumbents vs recent hires", "p25 → p74"), ("Cost to close the gap", "$342K"),
      ("Exposure carried", "$2.33M"), ("Payback", "no attrition yet")],
     "Invisible in the segment average, which is healthy. Long-tenured engineers sit at "
     "the 25th percentile of their own peer group while recent hires enter at the 74th. "
     "Attrition has not reacted — this is the only finding made before it did, and the "
     "only time an equity correction is cheap."),
    ("Singapore · Risk & Compliance IC2", "ORGANIZATIONAL", LIGHT,
     [("Compa-ratio", "1.01"), ("Attrition vs baseline", "2.07×"),
      ("Span of control", "15.1 vs 6.7"), ("Cost of a pay rise", "not the lever"),
      ("Exposure carried", "$2.18M"), ("A supervisory layer", "$0.5M – $1.2M / yr")],
     "Attrition twice the baseline and pay at market. One manager for every fifteen "
     "people against a company norm of one in seven. The lever is supervision, and it "
     "is not free: closing the span to 10 needs three more managers, to 7 needs eight."),
]
for head, tag, tag_color, rows, note in FINDINGS:
    s = slide()
    title(s, head)
    chip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(1.45),
                              Inches(2.3), Inches(0.36))
    chip.fill.solid()
    chip.fill.fore_color.rgb = rgb(tag_color)
    chip.line.fill.background()
    chip.shadow.inherit = False
    ctf = chip.text_frame
    ctf.margin_left = ctf.margin_right = 0
    cp = ctf.paragraphs[0]
    cp.alignment = PP_ALIGN.CENTER
    cr = cp.add_run()
    cr.text = tag
    cr.font.size = Pt(10)
    cr.font.bold = True
    cr.font.color.rgb = rgb(WHITE)
    cr.font.name = BODY_FONT
    for i, (label, value) in enumerate(rows):
        col, row = i % 3, i // 3
        x, y = 0.75 + col * 4.2, 2.2 + row * 1.35
        card(s, x, y, 3.9, 1.12)
        text(s, x + 0.28, y + 0.18, 3.34, 0.34,
             [(label, 10.5, MUTED, False, BODY_FONT, 0)])
        text(s, x + 0.28, y + 0.52, 3.34, 0.45,
             [(value, 19, NAVY, False, HEAD_FONT, 0)])
    card(s, 0.75, 5.05, W - 1.5, 1.25)
    text(s, 1.15, 5.32, W - 2.3, 0.95, [(note, 13, BODY, False, BODY_FONT, 0)])

# ---------------------------------------------------------------- 10 · the number
s = slide()
title(s, "The smallest number that fixes most of it")
card(s, 0.75, 2.15, 6.0, 4.15)
text(s, 1.05, 2.42, 5.4, 0.4,
     [("Where a blanket 6 % across LATAM goes", 12, NAVY, True, BODY_FONT, 0)])
chart(s, 0.95, 2.95, 5.6, 2.6, ["Already in band", "Below minimum"],
      [("Spend", [5.11, 0.27])], XL_CHART_TYPE.BAR_CLUSTERED, (OCHRE, NAVY), '[$-409]$#,##0.00"M"')
text(s, 1.05, 5.7, 5.4, 0.5,
     [("95 % of $5.38M reaches people the argument does not describe.",
       11.5, BODY, False, BODY_FONT, 0)])
card(s, 7.05, 2.15, 5.53, 4.15)
text(s, 7.35, 2.42, 5.0, 0.4, [("The comparison", 12, NAVY, True, BODY_FONT, 0)])
for i, (label, value, sub, color) in enumerate([
        ("Blanket 6 % across LATAM", "$5.38M", "23 people still below minimum after it", OCHRE),
        ("Targeted, two segments", "$1.15M", "covers $5.76M of exposure, 1 : 5", NAVY)]):
    y = 3.0 + i * 1.55
    text(s, 7.4, y, 4.9, 0.35, [(label, 11, MUTED, False, BODY_FONT, 0)])
    text(s, 7.4, y + 0.32, 4.9, 0.6, [(value, 30, color, False, HEAD_FONT, 0)])
    text(s, 7.4, y + 0.95, 4.9, 0.35, [(sub, 10.5, MUTED, False, BODY_FONT, 0)])

# ---------------------------------------------------------------- 11 · sensitivity
s = slide()
title(s, "The assumption it all rests on",
      "Every payback here depends on what replacing someone costs. That figure is "
      "declared, not measured — so here is the recommendation at the whole range.")
card(s, 0.75, 2.55, 6.0, 3.4)
text(s, 1.05, 2.8, 5.4, 0.4,
     [("Payback on the pay adjustment, months", 12, NAVY, True, BODY_FONT, 0)])
chart(s, 0.95, 3.25, 5.6, 2.45,
      ["1.0× salary", "1.5× salary", "2.0× salary"],
      [("Months", [7.3, 4.8, 3.6])], XL_CHART_TYPE.COLUMN_CLUSTERED,
      (LIGHT, BLUE, NAVY), '[$-409]#,##0.0')
card(s, 7.05, 2.55, 5.53, 3.4)
text(s, 7.4, 2.85, 4.9, 2.9,
     [("The spend does not move", 15, NAVY, True, BODY_FONT, 6),
      ("$1.15M at every assumption. The replacement factor prices the benefit, never "
       "the intervention — lifting a salary costs what it costs regardless of what "
       "replacing that person would have cost.", 12.5, BODY, False, BODY_FONT, 10),
      ("There is no breaking point in the range", 15, NAVY, True, BODY_FONT, 6),
      ("At the most conservative assumption anyone defends — 1.0× salary, the bottom of "
       "the published range — the recommendation still covers 3.3 times its cost and "
       "pays back inside eight months.", 12.5, BODY, False, BODY_FONT, 0)])
text(s, 0.75, 6.15, W - 1.5, 0.6,
     [("HR literature puts replacement cost between 0.5× and 2× annual salary, rising "
       "with specialization and ramp time. The model applies 0.5× to Operations through "
       "2.0× to managers.", 11, MUTED, False, BODY_FONT, 0)])

# ---------------------------------------------------------------- 12 · São Paulo
s = slide()
title(s, "The raise we did not recommend",
      "São Paulo Sales, 122 people, is the most under-paid population in the company "
      "on base salary. It is not on the action list.")
for i, (label, value, sub, color) in enumerate([
        ("Base compa-ratio", "0.88", "lowest of any segment", OCHRE),
        ("Total target cash", "0.99", "41 % variable against a 25 % norm", NAVY),
        ("Attrition vs baseline", "1.17×", "inside the noise, p = 0.49", NAVY)]):
    stat(s, 0.75 + i * 4.2, 2.65, 3.9, label, value, sub, color)
card(s, 0.75, 4.75, W - 1.5, 1.95)
text(s, 1.15, 5.05, W - 2.3, 1.4,
     [("Sales people in São Paulo are paid through variable compensation, and the "
       "market there expects it. Reading base salary alone would have recommended a "
       "raise for 122 people who were never underpaid — the single most expensive "
       "mistake available in this dataset, and the reason the rule reads total target "
       "cash rather than base.", 13, BODY, False, BODY_FONT, 0)])

# ---------------------------------------------------------------- 13 · risks
s = slide()
title(s, "What could make this wrong")
RISKS = [
    ("Attrition may not be caused by pay",
     "Bogotá shows low pay and high attrition together. That is a correlation. If people "
     "are leaving for managers, workload or career path, $811K buys nothing.",
     "Pilot it. Treat one half, hold the other."),
    ("Fixing IC3–IC4 alone inverts the ladder",
     "IC5 sits at 0.887 and IC6 at 0.928. Lift IC3–IC4 to 0.95 and the levels above them "
     "are paid less than the levels below.",
     "Extend to the whole family: $194K more."),
    ("Singapore's fix is priced at zero here",
     "A supervisory layer is a salary budget too. Three to eight managers, depending on "
     "the span targeted.",
     "Price it before committing: $0.5M – $1.2M."),
    ("The market survey is a point in time",
     "Every compa-ratio is measured against one survey year. If the market moved, the "
     "gap is measured against a stale midpoint.",
     "Re-run at the next survey before extending."),
]
for i, (head, rest, mit) in enumerate(RISKS):
    col, row = i % 2, i // 2
    x, y = 0.75 + col * 6.15, 2.15 + row * 2.35
    card(s, x, y, 5.85, 2.1)
    text(s, x + 0.3, y + 0.22, 5.25, 1.45,
         [(head, 14, NAVY, True, BODY_FONT, 5),
          (rest, 11.5, BODY, False, BODY_FONT, 0)])
    text(s, x + 0.3, y + 1.62, 5.25, 0.35,
         [("Mitigation:  ", 11, OCHRE, True, BODY_FONT, 0)])
    text(s, x + 1.15, y + 1.62, 4.4, 0.35,
         [(mit, 11, BODY, False, BODY_FONT, 0)])

# ---------------------------------------------------------------- 14 · the pilot
s = slide()
title(s, "The pilot the data already designed",
      "Bogotá Data & Analytics splits into two levels that are statistical twins. "
      "Treat one, hold the other — a controlled comparison, not a before-and-after.")
card(s, 0.75, 2.6, 5.85, 1.7)
text(s, 1.05, 2.85, 5.25, 1.25,
     [("IC3 — treated", 13, NAVY, True, BODY_FONT, 4),
      ("74 people · compa-ratio 0.853 · attrition 34.0 % · 24 exits\n"
       "Lift to 0.95 of market: $349K", 12, BODY, False, BODY_FONT, 0)])
card(s, 6.9, 2.6, 5.65, 1.7)
text(s, 7.2, 2.85, 5.05, 1.25,
     [("IC4 — control", 13, MUTED, True, BODY_FONT, 4),
      ("76 people · compa-ratio 0.851 · attrition 31.5 % · 25 exits\n"
       "No change until the pilot reads out", 12, BODY, False, BODY_FONT, 0)])
card(s, 0.75, 4.5, W - 1.5, 2.2)
text(s, 1.15, 4.75, 5.4, 1.7,
     [("When it can answer", 14, NAVY, True, BODY_FONT, 5),
      ("3 months  ·  p = 0.09  ·  cannot tell\n"
       "6 months  ·  p = 0.07  ·  suggestive\n"
       "9 months  ·  p = 0.01  ·  conclusive", 13, BODY, False, BODY_FONT, 0)])
text(s, 6.9, 4.75, 5.4, 1.7,
     [("But it can fail fast", 14, OCHRE, True, BODY_FONT, 5),
      ("Six or more exits from IC3 in the first three months would be unlikely if the "
       "adjustment were working. That is the stop-loss: killing a failure is quicker "
       "than proving a success.", 12.5, BODY, False, BODY_FONT, 0)])

# ---------------------------------------------------------------- 15 · next steps
s = slide()
title(s, "What happens next")
STEPS = [
    ("Now", "Approve $349K", "Bogotá DAT IC3 only. Dublin's $342K equity correction "
     "runs in parallel — it needs no pilot, the gap is measured, not inferred."),
    ("Month 1", "Adjustments effective", "Communicated as a market correction, not as a "
     "retention payment. IC4 unchanged and unaware."),
    ("Month 3", "Stop-loss review", "Six or more exits from IC3 and the pay hypothesis "
     "is wrong. Stop, and re-open the segment as an organizational question."),
    ("Month 9", "Read the pilot", "Conclusive at p = 0.01 if it worked. Extend to IC4 "
     "and the rest of the ladder: $462K plus $194K."),
]
for i, (when, what, detail) in enumerate(STEPS):
    x = 0.75 + i * 3.13
    card(s, x, 2.3, 2.88, 3.5)
    text(s, x + 0.26, 2.55, 2.36, 0.32, [(when, 11, OCHRE, True, BODY_FONT, 0)])
    text(s, x + 0.26, 2.92, 2.36, 0.85, [(what, 16, NAVY, True, HEAD_FONT, 0)])
    text(s, x + 0.26, 3.85, 2.36, 1.75, [(detail, 11, BODY, False, BODY_FONT, 0)])
card(s, 0.75, 6.05, W - 1.5, 0.85, fill=NAVY, line=None)
text(s, 1.15, 6.28, W - 2.3, 0.45,
     [("Total committed today: $691K. The remaining $656K is contingent on a result.",
       13, WHITE, True, BODY_FONT, 0)])

# ---------------------------------------------------------------- 16 · closing
s = slide(dark=True)
title(s, "What this proves, and what it does not", dark=True)
card(s, 0.75, 2.35, 5.85, 3.15, fill="0A2050", line="1B3466")
text(s, 1.1, 2.62, 5.2, 2.6,
     [("It was scored, not asserted", 15, WHITE, True, BODY_FONT, 7),
      ("The rule was run against an answer key written before the data existed: four "
       "anomalies built in, two of them designed to be dismissed.",
       12.5, PALE, False, BODY_FONT, 9),
      ("Precision 100 %  ·  Recall 100 %", 17, WHITE, True, HEAD_FONT, 9),
      ("Every figure is computed twice — once in DAX inside the model, once in Python "
       "from the warehouse — and the two are reconciled cell by cell. 77 automated "
       "tests. A wrong number fails a build, not a meeting.",
       12.5, PALE, False, BODY_FONT, 0)])
card(s, 6.9, 2.35, 5.65, 3.15, fill="0A2050", line="1B3466")
text(s, 7.25, 2.62, 5.05, 2.6,
     [("The limits, stated", 15, WHITE, True, BODY_FONT, 7),
      ("The thresholds are calibrated on this organization. A company with a different "
       "baseline needs them re-derived, not copied.", 12.5, PALE, False, BODY_FONT, 9),
      ("A perfect score on data built to be found is not a perfect score on real data. "
       "What this demonstrates is that the logic separates a real signal from two "
       "plausible decoys.", 12.5, PALE, False, BODY_FONT, 9),
      ("And it is why the recommendation is a pilot rather than a rollout.",
       12.5, WHITE, True, BODY_FONT, 0)])

prs.save(OUT)


def recolor_hyperlinks(path, colour):
    """Repaint the theme's hyperlink colour.

    A run carrying <a:hlinkClick> takes its colour from the theme, not from its own
    solidFill, so setting font.color on the run has no effect: the link renders in
    Office's default hyperlink blue, which on a navy title slide is invisible. The
    colour lives in the theme part and nowhere python-pptx exposes, so the file is
    repacked with the two values replaced.
    """
    import re
    import shutil
    import zipfile

    src = path.with_suffix(".tmp")
    shutil.move(path, src)
    with zipfile.ZipFile(src) as zin, zipfile.ZipFile(
            path, "w", zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            blob = zin.read(item.filename)
            if item.filename.startswith("ppt/theme/"):
                xml = blob.decode("utf-8")
                xml = re.sub(r'(<a:hlink><a:srgbClr val=")[0-9A-Fa-f]{6}',
                             r"\g<1>" + colour, xml)
                xml = re.sub(r'(<a:folHlink><a:srgbClr val=")[0-9A-Fa-f]{6}',
                             r"\g<1>" + colour, xml)
                blob = xml.encode("utf-8")
            zout.writestr(item, blob)
    src.unlink()


recolor_hyperlinks(OUT, LIGHT)
print(f"wrote {OUT.name} - {len(prs.slides.__iter__.__self__._sldIdLst)} slides, "
      f"{OUT.stat().st_size/1024:.0f} KB")
